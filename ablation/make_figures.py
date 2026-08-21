"""Generate editable PowerPoint versions of the architecture figures.

    python ablation/make_figures.py --config configs/full_legacy.yaml

Every box, channel count and grid size is derived from the model YAML rather than
drawn by hand, so the figure cannot drift from the network the way the current
Fig. 5 and Fig. 7 have. Shapes are real PowerPoint shapes and the arrows are
connected to them, so dragging a box keeps its arrows attached.

Slide 1 replaces Fig. 5, slide 2 replaces Fig. 7.
"""

from __future__ import annotations

import argparse
import math
import sys
from dataclasses import dataclass
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
INPUT_SIZE = 480

# Palette follows the existing figures so the replacement does not look foreign.
PALETTE: dict[str, tuple[str, str]] = {
    "Conv":                ("FCE4C8", "000000"),
    "C3k2":                ("BDD7EE", "000000"),
    "SPPF":                ("2E75B6", "FFFFFF"),
    "C2PSA":               ("FFC000", "000000"),
    "CrossAttentionBlock": ("E8503A", "FFFFFF"),
    "MultiScaleCBAM":      ("F2A9D2", "000000"),
    "MSCBAMFixed":         ("F2A9D2", "000000"),
    "StandardCBAM":        ("D9B3E6", "000000"),
    "nn.Upsample":         ("F8CBAD", "000000"),
    "Concat":              ("A9D18E", "000000"),
    "OBB":                 ("4472C4", "FFFFFF"),
    "Detect":              ("4472C4", "FFFFFF"),
}
DEFAULT_FILL = ("D9D9D9", "000000")


@dataclass(frozen=True)
class Layer:
    index: int
    module: str
    sources: list[int]
    channels: int
    grid: int
    stage: str


def scaled(c: int, width: float, max_channels: int) -> int:
    return math.ceil(min(c, max_channels) * width / 8) * 8


def trace(config: Path, scale: str = "s") -> list[Layer]:
    """Resolve every layer's channel count and spatial grid, as parse_model would."""
    spec = yaml.safe_load(config.read_text())
    width, max_channels = spec["scales"][scale][1], spec["scales"][scale][2]
    entries = spec["backbone"] + spec["head"]
    n_backbone = len(spec["backbone"])

    ch: dict[int, int] = {-1: 3}
    grid: dict[int, int] = {-1: INPUT_SIZE}
    layers: list[Layer] = []

    for i, entry in enumerate(entries):
        frm, module = entry[0], entry[2]
        args = entry[3] if len(entry) > 3 else []
        resolve = lambda x: x if x >= 0 else i + x
        sources = [resolve(x) for x in (frm if isinstance(frm, list) else [frm])]

        if module == "Concat":
            c, g = sum(ch[j] for j in sources), grid[sources[0]]
        elif module == "nn.Upsample":
            c, g = ch[i - 1], grid[i - 1] * 2
        elif module == "CrossAttentionBlock":
            c, g = scaled(args[0], width, max_channels), grid[sources[1]]
        elif module in ("OBB", "Detect"):
            c, g = 0, grid[sources[0]]
        elif module in ("MultiScaleCBAM", "MSCBAMFixed", "StandardCBAM"):
            c, g = ch[i - 1], grid[i - 1]
        else:
            c = scaled(args[0], width, max_channels)
            g = grid[i - 1] // 2 if len(args) > 2 and args[2] == 2 else grid[i - 1]

        ch[i], grid[i] = c, g
        layers.append(
            Layer(i, module, sources, c, g, "backbone" if i < n_backbone else "head")
        )
    return layers


def box(slide, x: Emu, y: Emu, w: Emu, h: Emu, text: str, module: str, pt: float = 8.5):
    fill, font = PALETTE.get(module, DEFAULT_FILL)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string("404040")
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False

    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Emu(18000)
    frame.margin_top = frame.margin_bottom = 0
    for n, line in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if n == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line
        run.font.size = Pt(pt if n == 0 else pt - 1.5)
        run.font.bold = n == 0
        run.font.color.rgb = RGBColor.from_string(font)
        run.font.name = "Calibri"
    return shape


def connect(slide, a, b, skip: bool = False):
    kind = MSO_CONNECTOR.ELBOW if skip else MSO_CONNECTOR.STRAIGHT
    c = slide.shapes.add_connector(kind, a.left, a.top, b.left, b.top)
    # 3 = right edge, 1 = left edge, 0 = top, 2 = bottom
    c.begin_connect(a, 2 if skip else 3)
    c.end_connect(b, 0 if skip else 1)
    c.line.color.rgb = RGBColor.from_string("7F9BB5" if skip else "31506B")
    c.line.width = Pt(1.75 if not skip else 1.0)
    return c


def title(slide, text: str, sub: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.16), Inches(12.6), Inches(0.62))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(15), True, "Calibri"
    p = tf.add_paragraph()
    r2 = p.add_run()
    r2.text = sub
    r2.font.size, r2.font.name = Pt(9.5), "Calibri"
    r2.font.color.rgb = RGBColor.from_string("595959")


def architecture_slide(prs, layers: list[Layer], config_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(
        slide,
        "Customised YOLOv11 architecture for text detection",
        f"Every block derived from {config_name}. Channels shown for the s scale "
        f"(width 0.50); grids for a {INPUT_SIZE}x{INPUT_SIZE} input. "
        "Grey elbow arrows are skip connections.",
    )

    cols, pad_x, top = 6, Inches(0.32), Inches(1.15)
    pitch_x = (SLIDE_W - 2 * pad_x) / cols
    bw, bh = pitch_x - Inches(0.17), Inches(0.78)
    rows = math.ceil(len(layers) / cols)
    pitch_y = (SLIDE_H - top - Inches(0.3)) / rows

    shapes: dict[int, object] = {}
    for layer in layers:
        row, col = divmod(layer.index, cols)
        x = pad_x + pitch_x * col + (pitch_x - bw) / 2
        y = top + pitch_y * row + (pitch_y - bh) / 2

        name = layer.module.replace("nn.", "")
        if layer.module in ("OBB", "Detect"):
            levels = "/".join(str(next(l for l in layers if l.index == s).grid)
                              for s in layer.sources)
            caption = f"{layer.index}  {name}\nP3/P4/P5  {levels}"
        else:
            caption = f"{layer.index}  {name}\nc={layer.channels}  {layer.grid}x{layer.grid}"
        shapes[layer.index] = box(slide, x, y, bw, bh, caption, layer.module)

    for layer in layers:
        for src in layer.sources:
            if src < 0:
                continue
            connect(slide, shapes[src], shapes[layer.index], skip=src != layer.index - 1)

    # Anchor each stage label over its own first block: with six columns the
    # backbone/neck boundary falls mid-row, so a row-start label would sit above
    # blocks belonging to the other stage.
    first_head = next(l.index for l in layers if l.stage == "head")
    for label, first in (("Backbone", 0), ("Neck and head", first_head)):
        row, col = divmod(first, cols)
        tb = slide.shapes.add_textbox(pad_x + pitch_x * col + (pitch_x - bw) / 2,
                                      top + pitch_y * row - Inches(0.21),
                                      Inches(2.4), Inches(0.22))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = label
        r.font.size, r.font.bold, r.font.name = Pt(9), True, "Calibri"
        r.font.color.rgb = RGBColor.from_string("595959")


def cross_attention_slide(prs, layers: list[Layer]) -> None:
    ca = next(l for l in layers if l.module == "CrossAttentionBlock")
    kv_src, q_src = ca.sources[0], ca.sources[1]
    kv = next(l for l in layers if l.index == kv_src)
    q = next(l for l in layers if l.index == q_src)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(
        slide,
        "Cross-Attention block",
        f"Query from layer {q.index} (P{int(math.log2(INPUT_SIZE // q.grid))}/"
        f"{INPUT_SIZE // q.grid}); key and value from layer {kv.index} "
        f"(P{int(math.log2(INPUT_SIZE // kv.grid))}/{INPUT_SIZE // kv.grid}). "
        "Shapes are for the s scale.",
    )

    y_q, y_kv, bh = Inches(1.75), Inches(4.35), Inches(1.05)
    stages = [
        (Inches(0.45), Inches(2.5), f"Layer {q.index} output\n(B, {q.channels}, {q.grid}, {q.grid})", "Conv", y_q),
        (Inches(0.45), Inches(2.5), f"Layer {kv.index} output\n(B, {kv.channels}, {kv.grid}, {kv.grid})", "C3k2", y_kv),
        (Inches(3.35), Inches(2.2), "Query Projection", "OBB", y_q),
        (Inches(3.35), Inches(2.2), f"Key & Value Projection\n1x1 conv, {kv.channels} to {ca.channels}", "CrossAttentionBlock", y_kv),
    ]
    made = [box(slide, x, y, w, bh, t, m, pt=10) for x, w, t, m, y in stages]

    attn = box(slide, Inches(6.1), Inches(2.85), Inches(2.6), Inches(1.5),
               f"Scaled Dot-Product\nCross-Attention\nh = 8 heads", "MultiScaleCBAM", pt=11)
    out = box(slide, Inches(9.15), Inches(2.85), Inches(2.0), Inches(1.5),
              "Output Projection\nAdd & LayerNorm", "Concat", pt=10)
    result = box(slide, Inches(11.5), Inches(2.85), Inches(1.55), Inches(1.5),
                 f"Output\n(B, {ca.channels}, {ca.grid}, {ca.grid})\nto BiFPN neck", "SPPF", pt=9.5)

    for a, b in ((made[0], made[2]), (made[1], made[3]), (made[2], attn),
                 (made[3], attn), (attn, out), (out, result)):
        connect(slide, a, b)

    note = slide.shapes.add_textbox(Inches(0.45), Inches(6.35), Inches(12.4), Inches(0.7))
    tf = note.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = (
        f"Query sequence length {q.grid * q.grid}; key/value sequence length "
        f"{kv.grid * kv.grid}. The key/value source is a plain convolutional feature "
        f"map, not an attention-refined one."
    )
    r.font.size, r.font.name = Pt(9.5), "Calibri"
    r.font.color.rgb = RGBColor.from_string("595959")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--config", type=Path,
                        default=Path(__file__).parent / "configs/full_legacy.yaml")
    parser.add_argument("--out", type=Path,
                        default=Path(__file__).parent / "figures_editable.pptx")
    parser.add_argument("--scale", default="s")
    opts = parser.parse_args()

    layers = trace(opts.config, opts.scale)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    architecture_slide(prs, layers, opts.config.name)
    cross_attention_slide(prs, layers)
    prs.save(opts.out)

    ca = next(l for l in layers if l.module == "CrossAttentionBlock")
    print(f"wrote {opts.out}  ({len(layers)} layers, scale {opts.scale})")
    print(f"  slide 1  architecture, {len(layers)} blocks")
    print(f"  slide 2  cross-attention: Q=L{ca.sources[1]}, K/V=L{ca.sources[0]}, "
          f"out c={ca.channels} {ca.grid}x{ca.grid}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
