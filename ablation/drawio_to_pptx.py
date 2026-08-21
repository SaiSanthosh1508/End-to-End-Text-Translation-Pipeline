"""Convert a .drawio diagram into PowerPoint shapes.

    python ablation/drawio_to_pptx.py architecture-final.drawio -o figure.pptx

draw.io has no PowerPoint export, and exporting an image gives a picture rather than
something editable. This reads the diagram XML and rebuilds it as native shapes, so
whatever was tweaked in draw.io arrives in PowerPoint still editable: boxes keep their
fill, outline and two-line label, and connectors are rebuilt segment by segment from
their waypoints so the routing is preserved exactly.

Connectors become straight segments joined at their corners rather than PowerPoint
elbow connectors, because an elbow connector re-routes itself and would discard the
lane routing the diagram depends on. Only the final segment carries the arrow head.
"""

from __future__ import annotations

import argparse
import re
import sys
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Pt

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
EMU_PER_IN = 914400
MARGIN_IN = 0.25


@dataclass
class Vertex:
    ident: str
    x: float
    y: float
    w: float
    h: float
    style: dict[str, str]
    raw_style: str
    lines: list[tuple[str, float, bool]] = field(default_factory=list)

    @property
    def is_label(self) -> bool:
        """A draw.io text shape: bare caption, no box drawn around it."""
        return self.raw_style.startswith("text;")


@dataclass
class Edge:
    source: str
    target: str
    style: dict[str, str]
    points: list[tuple[float, float]]


def parse_style(raw: str) -> dict[str, str]:
    out: dict[str, str] = {}
    for part in (raw or "").split(";"):
        key, sep, value = part.partition("=")
        if sep:
            out[key.strip()] = value.strip()
    return out


def parse_label(value: str) -> list[tuple[str, float, bool]]:
    """Split a draw.io HTML label into (text, size_px, bold) lines."""
    if not value:
        return []
    lines: list[tuple[str, float, bool]] = []
    for chunk in re.split(r"<br\s*/?>", value):
        size_match = re.search(r"font-size:\s*(\d+(?:\.\d+)?)px", chunk)
        bold = "<b" in chunk or "<strong" in chunk
        text = re.sub(r"<[^>]+>", "", chunk)
        text = (
            text.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
            .replace("&quot;", '"').replace("&apos;", "'").replace("&nbsp;", " ")
        ).strip()
        if text:
            lines.append((text, float(size_match.group(1)) if size_match else 0.0, bold))
    return lines


def read(path: Path) -> tuple[dict[str, Vertex], list[Edge]]:
    root = ET.parse(path).getroot()
    if root.tag != "mxGraphModel":
        found = root.find(".//mxGraphModel")
        if found is None:
            raise SystemExit(f"{path} has no mxGraphModel; is it a compressed .drawio?")
        root = found

    vertices: dict[str, Vertex] = {}
    edges: list[Edge] = []
    for cell in root.findall(".//mxCell"):
        geometry = cell.find("mxGeometry")
        style = parse_style(cell.get("style", ""))
        if cell.get("vertex") == "1" and geometry is not None:
            vertices[cell.get("id")] = Vertex(
                cell.get("id"),
                float(geometry.get("x", 0)), float(geometry.get("y", 0)),
                float(geometry.get("width", 0)), float(geometry.get("height", 0)),
                style, cell.get("style", ""), parse_label(cell.get("value", "")),
            )
        elif cell.get("edge") == "1":
            points = [
                (float(p.get("x")), float(p.get("y")))
                for array in (geometry.findall("Array") if geometry is not None else [])
                if array.get("as") == "points"
                for p in array.findall("mxPoint")
            ]
            edges.append(Edge(cell.get("source"), cell.get("target"), style, points))
    return vertices, edges


def anchor(v: Vertex, style: dict[str, str], prefix: str, other: Vertex) -> tuple[float, float]:
    """Where a connector meets a block: the pinned side if the style names one."""
    if f"{prefix}X" in style:
        return (v.x + float(style[f"{prefix}X"]) * v.w,
                v.y + float(style[f"{prefix}Y"]) * v.h)
    dx, dy = (other.x + other.w / 2) - (v.x + v.w / 2), (other.y + other.h / 2) - (v.y + v.h / 2)
    if abs(dx) > abs(dy):
        return (v.x + v.w if dx > 0 else v.x, v.y + v.h / 2)
    return (v.x + v.w / 2, v.y + v.h if dy > 0 else v.y)


def orthogonal(points: list[tuple[float, float]], start_horizontal: bool) -> list[tuple[float, float]]:
    out = [points[0]]
    horizontal = start_horizontal
    for nxt in points[1:]:
        x0, y0 = out[-1]
        x1, y1 = nxt
        if abs(x0 - x1) > 0.5 and abs(y0 - y1) > 0.5:
            out.append((x1, y0) if horizontal else (x0, y1))
        out.append(nxt)
        if abs(out[-1][0] - out[-2][0]) > 0.5:
            horizontal = False
        elif abs(out[-1][1] - out[-2][1]) > 0.5:
            horizontal = True
    return out


def route(edge: Edge, vertices: dict[str, Vertex]) -> list[tuple[float, float]]:
    a, b = vertices[edge.source], vertices[edge.target]
    start = anchor(a, edge.style, "exit", b)
    end = anchor(b, edge.style, "entry", a)
    horizontal = abs(start[0] - (a.x + a.w / 2)) > abs(start[1] - (a.y + a.h / 2))
    return orthogonal([start, *edge.points, end], horizontal)


def add_box(slide, v: Vertex, to_emu, scale: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, *to_emu(v.x, v.y), Emu(int(v.w * scale)), Emu(int(v.h * scale))
    )
    shape.adjustments[0] = 0.16
    if v.is_label:
        shape.fill.background()
        shape.line.fill.background()
    else:
        fill = v.style.get("fillColor", "#FFFFFF")
        if fill.lower() == "none":
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip("#"))
        shape.line.color.rgb = RGBColor.from_string(
            v.style.get("strokeColor", "#000000").lstrip("#")
        )
        shape.line.width = Pt(float(v.style.get("strokeWidth", 1)))
        if v.style.get("dashed") == "1":
            shape.line.dash_style = 4  # MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False

    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Emu(27000)
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = (
        MSO_ANCHOR.TOP if v.style.get("verticalAlign") == "top" else MSO_ANCHOR.MIDDLE
    )
    base = float(v.style.get("fontSize", 12))
    colour = RGBColor.from_string(v.style.get("fontColor", "#000000").lstrip("#"))
    for n, (text, size, bold) in enumerate(v.lines):
        para = frame.paragraphs[0] if n == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.LEFT if v.style.get("align") == "left" else PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt((size or base) * 0.75)   # draw.io px -> points
        run.font.bold = bold or v.style.get("fontStyle") in ("1", "3")
        run.font.italic = v.style.get("fontStyle") in ("2", "3")
        run.font.color.rgb = colour
        run.font.name = "Helvetica"


def add_edge(slide, points: list[tuple[float, float]], style: dict[str, str], to_emu) -> None:
    colour = RGBColor.from_string(style.get("strokeColor", "#000000").lstrip("#"))
    width = Pt(float(style.get("strokeWidth", 1)))
    for n, (start, end) in enumerate(zip(points, points[1:])):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *to_emu(*start), *to_emu(*end))
        line.line.color.rgb = colour
        line.line.width = width
        if n == len(points) - 2:
            # python-pptx exposes no arrowhead API, and only the last segment should
            # carry one, so the DrawingML element is appended directly.
            line.line._get_or_add_ln().append(
                parse_xml(
                    '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/'
                    '2006/main" type="triangle" w="med" len="med"/>'
                )
            )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source", type=Path, help="the .drawio file to convert")
    parser.add_argument("-o", "--out", type=Path, help="output .pptx (default: alongside source)")
    opts = parser.parse_args()

    vertices, edges = read(opts.source)
    if not vertices:
        raise SystemExit(f"no shapes found in {opts.source}")

    span_x = max(v.x + v.w for v in vertices.values())
    span_y = max(v.y + v.h for v in vertices.values())
    for edge in edges:
        for px, py in edge.points:
            span_x, span_y = max(span_x, px), max(span_y, py)
    origin_x = min([v.x for v in vertices.values()] + [p[0] for e in edges for p in e.points])
    origin_y = min([v.y for v in vertices.values()] + [p[1] for e in edges for p in e.points])

    scale = min(
        (SLIDE_W_IN - 2 * MARGIN_IN) * EMU_PER_IN / (span_x - origin_x),
        (SLIDE_H_IN - 2 * MARGIN_IN) * EMU_PER_IN / (span_y - origin_y),
    )
    offset_x = (SLIDE_W_IN * EMU_PER_IN - (span_x - origin_x) * scale) / 2
    offset_y = (SLIDE_H_IN * EMU_PER_IN - (span_y - origin_y) * scale) / 2

    def to_emu(x: float, y: float) -> tuple[Emu, Emu]:
        return Emu(int((x - origin_x) * scale + offset_x)), Emu(int((y - origin_y) * scale + offset_y))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(int(SLIDE_W_IN * EMU_PER_IN)), Emu(int(SLIDE_H_IN * EMU_PER_IN))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    labels = [v for v in vertices.values() if v.is_label]
    panels = [v for v in vertices.values()
              if not v.is_label and v.w * v.h > 40000]
    blocks = [v for v in vertices.values() if v not in panels and v not in labels]
    for panel in panels:                       # panels first so blocks sit above them
        add_box(slide, panel, to_emu, scale)
    for edge in edges:
        add_edge(slide, route(edge, vertices), edge.style, to_emu)
    for block in blocks + labels:
        add_box(slide, block, to_emu, scale)

    out = opts.out or opts.source.with_suffix(".pptx")
    prs.save(out)
    print(f"wrote {out}")
    print(f"  {len(blocks)} blocks, {len(panels)} panels, {len(labels)} captions, "
          f"{len(edges)} connectors")
    print(f"  {span_x - origin_x:.0f}x{span_y - origin_y:.0f} diagram units -> "
          f"{(span_x - origin_x) * scale / EMU_PER_IN:.2f}x"
          f"{(span_y - origin_y) * scale / EMU_PER_IN:.2f} in")
    return 0


if __name__ == "__main__":
    sys.exit(main())
